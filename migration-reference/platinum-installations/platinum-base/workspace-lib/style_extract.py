"""
Style extraction helpers.

Extract design specs from reference files (images, PPTX, PDF) to drive
consistent report generation. Works with ReportTemplate for template-based
PowerPoint creation.

Usage:
    from style_extract import extract_from_pptx, extract_colors_from_image, build_style_spec, spec_to_css

    design = extract_from_pptx('/workspace/uploads/reference.pptx')
    colors = extract_colors_from_image('/workspace/uploads/screenshot.png')
    style = build_style_spec(design)
"""

import json
import logging
import os
import subprocess
from typing import Optional

from lxml import etree

logger = logging.getLogger('style_extract')


# XML namespace map for OOXML theme parsing
_NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

# Theme color slot names in canonical order
_COLOR_SLOTS = [
    "dk1", "lt1", "dk2", "lt2",
    "accent1", "accent2", "accent3", "accent4",
    "accent5", "accent6", "hlink", "folHlink",
]


def extract_colors_from_image(image_path, n=8):
    """Extract dominant colors from an image via palette quantization.

    Args:
        image_path: Path to an image file (PNG, JPG, etc.).
        n: Number of dominant colors to extract (default 8).

    Returns:
        List of hex color strings sorted by frequency (most frequent first).
        Returns empty list if Pillow is not available or image cannot be read.
    """
    try:
        from PIL import Image
    except ImportError:
        return []

    if not os.path.isfile(image_path):
        return []

    try:
        img = Image.open(image_path).convert("RGB")
        img = img.resize((150, 150), Image.LANCZOS)
        quantized = img.quantize(colors=n, method=Image.Quantize.MEDIANCUT)

        # Get palette (flat list of R, G, B values)
        palette = quantized.getpalette()
        if not palette:
            return []

        # Count pixel frequency per palette index
        pixel_counts = {}
        for pixel in quantized.getdata():
            pixel_counts[pixel] = pixel_counts.get(pixel, 0) + 1

        # Build (hex_color, count) pairs
        color_freq = []
        for idx, count in pixel_counts.items():
            if idx < n:
                r = palette[idx * 3]
                g = palette[idx * 3 + 1]
                b = palette[idx * 3 + 2]
                hex_color = f"#{r:02x}{g:02x}{b:02x}"
                color_freq.append((hex_color, count))

        # Sort by frequency descending
        color_freq.sort(key=lambda x: x[1], reverse=True)
        return [c[0] for c in color_freq]

    except Exception:
        return []


def extract_from_pptx(pptx_path):
    """Extract complete design spec from a PPTX file.

    Parses theme XML for colors and fonts, enumerates layouts with
    placeholders, extracts custom shape fill colors from content slides,
    and reads slide dimensions.

    Args:
        pptx_path: Path to a .pptx file.

    Returns:
        Dict with keys:
            - colors: dict of theme color slots (dk1, lt1, accent1, etc.)
            - shape_colors: list of hex colors found in slide shape fills
            - fonts: dict with 'heading' and 'body' font names
            - layouts: list of layout dicts with name and placeholders
            - slide_size: dict with 'width' and 'height' in inches
        Returns None if file cannot be read.
    """
    try:
        from pptx import Presentation
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.util import Emu
    except ImportError:
        return None

    if not os.path.isfile(pptx_path):
        return None

    try:
        prs = Presentation(pptx_path)
    except Exception:
        return None

    colors = {}
    fonts = {"heading": "Calibri", "body": "Calibri"}
    layouts = []
    shape_colors = []

    # --- Theme colors and fonts ---
    try:
        theme_part = prs.slide_masters[0].part.part_related_by(RT.THEME)
        theme_xml = theme_part.blob
        root = etree.fromstring(theme_xml)

        # Color scheme
        clr_scheme = root.find(".//a:clrScheme", _NSMAP)
        if clr_scheme is not None:
            for name in _COLOR_SLOTS:
                el = clr_scheme.find(f"a:{name}", _NSMAP)
                if el is not None:
                    srgb = el.find("a:srgbClr", _NSMAP)
                    sys_clr = el.find("a:sysClr", _NSMAP)
                    if srgb is not None:
                        colors[name] = "#" + srgb.get("val", "000000")
                    elif sys_clr is not None:
                        colors[name] = "#" + sys_clr.get("lastClr", "000000")

        # Font scheme
        font_scheme = root.find(".//a:fontScheme", _NSMAP)
        if font_scheme is not None:
            major = font_scheme.find("a:majorFont/a:latin", _NSMAP)
            minor = font_scheme.find("a:minorFont/a:latin", _NSMAP)
            if major is not None:
                fonts["heading"] = major.get("typeface", "Calibri")
            if minor is not None:
                fonts["body"] = minor.get("typeface", "Calibri")
    except Exception:
        pass

    # --- Layouts ---
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            ph_info = []
            for ph in layout.placeholders:
                ph_info.append({
                    "idx": ph.placeholder_format.idx,
                    "name": ph.name,
                    "type": str(ph.placeholder_format.type),
                })
            layouts.append({
                "name": layout.name,
                "placeholders": ph_info,
            })

    # --- Shape fill colors from content slides ---
    seen_colors = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                fill = shape.fill
                if fill.type is not None:
                    fc = fill.fore_color
                    if fc and fc.type is not None:
                        rgb = fc.rgb
                        if rgb:
                            hex_str = f"#{rgb}"
                            if hex_str.lower() not in seen_colors:
                                seen_colors.add(hex_str.lower())
                                shape_colors.append(hex_str)
            except Exception:
                continue

    # --- Slide dimensions ---
    width_in = round(prs.slide_width / 914400, 3)
    height_in = round(prs.slide_height / 914400, 3)

    return {
        "colors": colors,
        "shape_colors": shape_colors,
        "fonts": fonts,
        "layouts": layouts,
        "slide_size": {"width": width_in, "height": height_in},
    }


def extract_from_pdf(pdf_path, page=0):
    """Extract a page from a PDF as a PIL Image.

    Tries pypdfium2 first (best quality), falls back to pypdf for
    embedded image extraction.

    Args:
        pdf_path: Path to a PDF file.
        page: Zero-based page index to extract (default 0).

    Returns:
        PIL Image of the rendered page, or None if extraction fails.
    """
    if not os.path.isfile(pdf_path):
        return None

    # Try pypdfium2 first
    try:
        import pypdfium2 as pdfium
        pdf = pdfium.PdfDocument(pdf_path)
        if page < len(pdf):
            pdfium_page = pdf[page]
            bitmap = pdfium_page.render(scale=2)
            pil_image = bitmap.to_pil()
            pdf.close()
            return pil_image
        pdf.close()
    except ImportError:
        pass
    except Exception:
        pass

    # Fallback: pypdf embedded images
    try:
        from pypdf import PdfReader
        from PIL import Image
        import io

        reader = PdfReader(pdf_path)
        if page < len(reader.pages):
            pdf_page = reader.pages[page]
            if "/XObject" in pdf_page["/Resources"]:
                xobjects = pdf_page["/Resources"]["/XObject"].get_object()
                for obj_name in xobjects:
                    xobj = xobjects[obj_name].get_object()
                    if xobj["/Subtype"] == "/Image":
                        data = xobj.get_data()
                        return Image.open(io.BytesIO(data))
    except ImportError:
        pass
    except Exception:
        pass

    return None


def pptx_to_images(pptx_path, output_dir=None):
    """Convert a PPTX file to PNG slide images.

    Uses LibreOffice headless to convert PPTX to PDF, then pypdfium2 to
    render each page as a high-resolution PNG.

    Args:
        pptx_path: Path to a .pptx file.
        output_dir: Directory to write slide PNGs (created if needed).
                    Defaults to /workspace/slides/<pptx_basename>/ to avoid
                    collisions when multiple PPTXs are processed.

    Returns:
        List of output PNG file paths, or empty list on failure.
    """
    if not os.path.isfile(pptx_path):
        return []

    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    if output_dir is None:
        output_dir = f"/workspace/slides/{base_name}"

    os.makedirs(output_dir, exist_ok=True)

    # Step 1: PPTX -> PDF via LibreOffice
    try:
        result = subprocess.run(
            [
                "soffice", "--headless", "--convert-to", "pdf",
                "--outdir", "/tmp", pptx_path,
            ],
            capture_output=True,
            timeout=120,
        )
        if result.returncode != 0:
            logger.warning(
                'LibreOffice conversion failed (rc=%d) for %s\nstdout: %s\nstderr: %s',
                result.returncode, pptx_path,
                result.stdout[:500] if result.stdout else '',
                result.stderr[:500] if result.stderr else '',
            )
            return []
    except (FileNotFoundError, subprocess.TimeoutExpired) as exc:
        logger.warning('LibreOffice not available or timed out for %s: %s', pptx_path, exc)
        return []
    pdf_path = f"/tmp/{base_name}.pdf"

    if not os.path.isfile(pdf_path):
        logger.warning('Expected PDF not found at %s after LibreOffice conversion of %s', pdf_path, pptx_path)
        return []

    # Step 2: PDF -> PNGs via pypdfium2
    try:
        import pypdfium2 as pdfium
    except ImportError:
        logger.warning('pypdfium2 not available — cannot render PDF to images')
        return []

    output_paths = []
    try:
        pdf = pdfium.PdfDocument(pdf_path)
        for i in range(len(pdf)):
            page = pdf[i]
            # scale=1.6 gives ~1536x864 for standard slides (13.33"x7.5" at 72dpi).
            # This fits under Claude's 1568px vision threshold while staying sharp.
            bitmap = page.render(scale=1.6)
            pil_image = bitmap.to_pil()
            out_path = os.path.join(output_dir, f"slide_{i + 1:03d}.png")
            pil_image.save(out_path, "PNG", optimize=True)
            output_paths.append(out_path)
        pdf.close()
    except Exception as exc:
        logger.warning('PDF rendering failed for %s: %s', pdf_path, exc)

    # Clean up temp PDF
    try:
        os.remove(pdf_path)
    except OSError:
        pass

    return output_paths


def build_style_spec(design_spec, vision_notes=""):
    """Build a unified style specification from extracted design data.

    Combines programmatic theme extraction with optional AI vision analysis
    notes into a single style spec suitable for driving report generation.

    Args:
        design_spec: Dict from extract_from_pptx() containing colors, fonts,
            and layouts.
        vision_notes: Optional free-text notes from AI vision analysis of
            slide screenshots (layout observations, styling details, etc.).

    Returns:
        Dict with keys:
            - palette: derived color palette (primary, accent, background,
              text, series)
            - fonts: heading and body font names
            - layout: layout mode ('dark' or 'light')
            - vision_notes: passthrough of vision analysis notes
            - raw_theme_colors: original theme color dict
    """
    if design_spec is None:
        design_spec = {}

    colors = design_spec.get("colors", {})

    # Build accent series from theme accents
    accent_keys = ["accent1", "accent2", "accent3", "accent4", "accent5", "accent6"]
    series = [colors[k] for k in accent_keys if k in colors]

    # Determine primary and accent
    primary = colors.get("accent1", "#3b82f6")
    accent = colors.get("accent2", "#ef4444")

    # Background and text from dk1/lt1
    bg = colors.get("dk1", "#0f172a")
    text = colors.get("lt1", "#f1f5f9")

    # Detect light vs dark mode
    mode = "dark" if _is_dark(bg) else "light"

    # If mode is light, swap bg and text semantics
    if mode == "light":
        bg = colors.get("lt1", "#f1f5f9")
        text = colors.get("dk1", "#0f172a")

    fonts = design_spec.get("fonts", {"heading": "Calibri", "body": "Calibri"})

    return {
        "palette": {
            "primary": primary,
            "accent": accent,
            "background": bg,
            "text": text,
            "series": series if series else [
                "#3b82f6", "#ef4444", "#22c55e",
                "#f59e0b", "#8b5cf6", "#ec4899",
            ],
        },
        "fonts": fonts,
        "layout": mode,
        "vision_notes": vision_notes,
        "raw_theme_colors": colors,
    }


def spec_to_css(design_spec, output_path=None):
    """Convert extracted PPTX design spec to CSS theme variables.

    Generates a :root block with CSS custom properties derived from the
    design spec's color scheme and font names. The variable names are
    stable across themes so HTML templates can reference them without
    knowing which brand is active.

    Args:
        design_spec: Output from extract_from_pptx(). Must contain at
            minimum a 'colors' dict and 'fonts' dict. Missing keys are
            filled with sensible defaults.
        output_path: If provided, writes the CSS string to this file
            path (parent directories are created automatically).

    Returns:
        CSS string containing :root { ... } with theme variables.
    """
    if design_spec is None:
        design_spec = {}

    colors = design_spec.get("colors", {})
    fonts = design_spec.get("fonts", {})

    # Determine light vs dark mode from dk1/lt1
    dk1 = colors.get("dk1", "#333333")
    lt1 = colors.get("lt1", "#ffffff")
    mode = "dark" if _is_dark(dk1) is False else "light"

    # In a typical PPTX theme dk1 is the dark text color and lt1 is
    # the light background.  Map to semantic roles accordingly.
    if mode == "light":
        bg = lt1
        text = dk1
    else:
        bg = dk1
        text = lt1

    primary = colors.get("accent1", "#0072CE")

    # text-light: a muted version of the text color. If dk2 is
    # available use it, otherwise derive from text.
    text_light = colors.get("dk2", "")
    if not text_light:
        text_light = _muted_hex(text) if text else "#666666"

    heading_font = fonts.get("heading", "Arial")
    body_font = fonts.get("body", "Arial")

    accents = []
    for i in range(1, 7):
        accents.append(colors.get(f"accent{i}", ""))

    sidebar = colors.get("accent1", primary)

    lines = [":root {"]
    lines.append(f"  --primary: {primary};")
    lines.append(f"  --bg: {bg};")
    lines.append(f"  --text: {text};")
    lines.append(f"  --text-light: {text_light};")
    lines.append(f"  --font-heading: {heading_font};")
    lines.append(f"  --font-body: {body_font};")
    for i, accent in enumerate(accents, 1):
        val = accent if accent else primary
        lines.append(f"  --accent{i}: {val};")
    lines.append(f"  --sidebar-color: {sidebar};")
    lines.append("}")
    lines.append("")  # trailing newline

    css = "\n".join(lines)

    if output_path:
        out_dir = os.path.dirname(output_path)
        if out_dir:
            os.makedirs(out_dir, exist_ok=True)
        with open(output_path, "w") as f:
            f.write(css)
        logger.info("Wrote theme CSS to %s", output_path)

    return css


def _muted_hex(hex_color):
    """Return a muted (lighter/darker) variant of a hex color.

    Moves the color 40% toward middle grey (#808080). Useful for
    deriving --text-light from --text when no dk2 is available.

    Args:
        hex_color: Color string like '#333333'.

    Returns:
        Hex color string.
    """
    h = hex_color.lstrip("#")
    if len(h) < 6:
        return "#666666"
    r = int(h[0:2], 16)
    g = int(h[2:4], 16)
    b = int(h[4:6], 16)
    # Blend 40% toward 0x80
    r = int(r + 0.4 * (0x80 - r))
    g = int(g + 0.4 * (0x80 - g))
    b = int(b + 0.4 * (0x80 - b))
    return f"#{r:02x}{g:02x}{b:02x}"


def _is_dark(hex_color):
    """Determine if a hex color is dark based on luminance.

    Uses the ITU-R BT.601 luminance formula.

    Args:
        hex_color: Color string like '#0f172a' or '0f172a'.

    Returns:
        True if the color's luminance is below 0.5 (dark).
    """
    h = hex_color.lstrip("#")
    if len(h) < 6:
        return True
    r = int(h[0:2], 16)
    g = int(h[2:4], 16)
    b = int(h[4:6], 16)
    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
    return luminance < 0.5
