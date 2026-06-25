#!/usr/bin/env python3
"""Deterministic PPTX report generator.

Takes a slide manifest + client PPTX template and produces a branded report
with all formatting rules baked in. The agent's job is to build the manifest
and gather data — this script handles layout, charts, and formatting.

What this handles (so the agent doesn't have to):
    - Adaptive chart layout based on title/insight/footer length
    - reverse_order on ALL bar charts (highest at top)
    - Legend hidden on single-series charts
    - Insight positioned below chart (never beside)
    - Number format '0"%"' (never '0%')
    - Empty data detection (skip chart, keep slide)
    - Adaptive font sizing for dense charts
    - Auto-fit text flags for PowerPoint rendering
    - Section divider gradient cycling
    - Placeholder access by idx (never by name)
    - Row filtering (Average, Mean, Total)

Usage:
    python generate_pptx.py --template template.pptx --manifest manifest.json --output report.pptx
    python generate_pptx.py --template t.pptx --manifest m.json --output r.pptx \\
        --colors '#99FFED,#89CDFF,#CA88FF,#FF99DD,#FFB288,#FFE588' --single-color '#AAFF89'
    python generate_pptx.py --validate-only existing.pptx
    python generate_pptx.py --template t.pptx --manifest m.json --output r.pptx --validate --thumbnails

Manifest format (JSON array of slides, or object with "config" + "slides"):

    [
      {"type": "title", "title": "Report Title"},
      {"type": "text", "title": "About", "body": "Body text..."},
      {"type": "section", "title": "Section Name"},
      {"type": "chart", "title": "Chart Title", "chart_type": "bar",
       "data_file": "table.json", "insight": "AI insight", "footer": "Source"},
      {"type": "thankyou", "title": "Thank You"}
    ]

Optional config block (CLI flags override these):

    {
      "config": {
        "colors": ["#99FFED", ...],
        "single_color": "#AAFF89",
        "sidebar_width": 1.5,
        "layout_map": {"title": "Title slide", ...},
        "placeholder_map": {"body": 17, ...}
      },
      "slides": [...]
    }
"""

import argparse
import json
import os
import sys

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Data loading — to_dataframe from the Platinum lib
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), '..'))
try:
    from platinum import to_dataframe
except ImportError:
    def to_dataframe(path):
        """Fallback: load PlatinumData JSON into a DataFrame."""
        with open(path) as f:
            data = json.load(f)
        rows = data.get('rows', [])
        if not rows:
            return pd.DataFrame()
        cols = [c.get('label', f'col{i}') for i, c in enumerate(data.get('columns', []))]
        records = {}
        for row in rows:
            label = row.get('label', '')
            vals = []
            for cell in row.get('cells', []):
                v = cell.get('colpc', cell.get('value', 0))
                vals.append(float(v) if v is not None else 0.0)
            records[label] = vals
        return pd.DataFrame.from_dict(records, orient='index', columns=cols[:len(next(iter(records.values())))] if records else cols)


# =============================================================================
# Constants
# =============================================================================

SLIDE_WIDTH = 13.33
SLIDE_HEIGHT = 7.50
DEFAULT_SIDEBAR_WIDTH = 1.5
BOTTOM_MARGIN = 0.2
TITLE_TOP = 0.25

DEFAULT_COLORS = [
    '#4472C4', '#ED7D31', '#A5A5A5', '#FFC000', '#5B9BD5', '#70AD47'
]

DEFAULT_EXCLUDE_ROWS = {'average', 'mean', 'avg', 'total'}

CHART_TYPE_MAP = {
    'bar': XL_CHART_TYPE.BAR_CLUSTERED,
    'clustered bar': XL_CHART_TYPE.BAR_CLUSTERED,
    'clustered_bar': XL_CHART_TYPE.BAR_CLUSTERED,
    'stacked bar': XL_CHART_TYPE.BAR_STACKED_100,
    'stacked_bar': XL_CHART_TYPE.BAR_STACKED_100,
    'stacked bar 100': XL_CHART_TYPE.BAR_STACKED_100,
    'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
    'clustered column': XL_CHART_TYPE.COLUMN_CLUSTERED,
    'clustered_column': XL_CHART_TYPE.COLUMN_CLUSTERED,
    'stacked column': XL_CHART_TYPE.COLUMN_STACKED_100,
    'stacked_column': XL_CHART_TYPE.COLUMN_STACKED_100,
    'stacked column 100': XL_CHART_TYPE.COLUMN_STACKED_100,
    'pie': XL_CHART_TYPE.PIE,
}

SLIDE_TYPE_ALIASES = {
    'title': 'title', 'title slide': 'title', 'first slide': 'title',
    'text': 'text', 'text slide': 'text',
    'summary': 'text', 'summary slide': 'text',
    'about': 'text', 'methodology': 'text',
    'chart': 'chart', 'graph': 'chart', 'graph slide': 'chart',
    'section': 'section', 'section slide': 'section',
    'break': 'section', 'break slide': 'section', 'divider': 'section',
    'appendix': 'appendix', 'appendix slide': 'appendix',
    'thankyou': 'thankyou', 'thank you': 'thankyou',
    'thank you slide': 'thankyou', 'end': 'thankyou',
    'quote': 'quote', 'quote slide': 'quote',
    'two column': 'two_column', 'two_column': 'two_column',
}

DEFAULT_LAYOUT_MAP = {
    'title': 'Title slide',
    'text': 'One column slide',
    'chart': 'Graph slide',
    'section': [
        'Section title slide gradient A',
        'Section title slide gradient B',
        'Section title slide gradient C',
        'Section title slide gradient D',
        'Section title slide gradient E',
    ],
    'appendix': 'Appendix slide',
    'thankyou': 'Thank you slide',
    'quote': 'Quote slide',
    'two_column': 'Two column slide',
}

DEFAULT_PLACEHOLDER_MAP = {
    'body': 17,
    'graph_footer': 22,
    'text_footer': 18,
    'section_footer': 18,
    'appendix_footer': 19,
}


# =============================================================================
# Text measurement
# =============================================================================

def measure_text_lines(text, font_size_pt, available_width_inches, font_name='Arial'):
    """Measure how many lines text will wrap to.

    Uses fonttools for real character widths when available.
    Falls back to character-count estimation.
    """
    if not text:
        return 1
    try:
        from fontTools.ttLib import TTFont

        font_paths = [
            '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
            '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
            '/usr/share/fonts/TTF/DejaVuSans.ttf',
        ]
        font_path = next((fp for fp in font_paths if os.path.exists(fp)), None)
        if not font_path:
            raise FileNotFoundError('No system font found')

        font = TTFont(font_path)
        cmap = font.getBestCmap()
        units_per_em = font['head'].unitsPerEm
        hmtx = font['hmtx']

        available_width_pt = available_width_inches * 72
        available_width_units = (available_width_pt / font_size_pt) * units_per_em
        space_width = hmtx['space'][0] if 'space' in hmtx.metrics else int(units_per_em * 0.25)

        lines = 1
        current_width = 0
        for word in text.split():
            word_width = space_width
            for ch in word:
                glyph_name = cmap.get(ord(ch))
                if glyph_name and glyph_name in hmtx.metrics:
                    word_width += hmtx[glyph_name][0]
                else:
                    word_width += int(units_per_em * 0.6)
            if current_width + word_width > available_width_units and current_width > 0:
                lines += 1
                current_width = word_width
            else:
                current_width += word_width
        return lines

    except Exception:
        chars_per_line = int(available_width_inches * 72 / font_size_pt * 1.8)
        if chars_per_line <= 0:
            return 1
        return max(1, -(-len(text) // chars_per_line))


# =============================================================================
# Layout calculators
# =============================================================================

def calculate_layout(title_text, sidebar_width=DEFAULT_SIDEBAR_WIDTH,
                     has_insight=False, insight_text='', has_footer=False):
    """Calculate chart position and title font size based on actual content.

    Returns: (left, top, width, height, title_font_pt)
             left/top/width/height as pptx EMU values (from Inches()),
             title_font_pt as an integer.
    """
    content_width = SLIDE_WIDTH - sidebar_width - 0.3
    title_area_width = content_width - 0.5

    title_font_pt = 24
    title_lines = measure_text_lines(title_text, 24, title_area_width)

    if title_lines >= 3:
        title_font_pt = 20
        title_lines = measure_text_lines(title_text, 20, title_area_width)

    if title_lines >= 4:
        title_font_pt = 18
        title_lines = measure_text_lines(title_text, 18, title_area_width)

    line_height = title_font_pt / 72
    title_height = title_lines * line_height * 1.4  # generous line spacing
    chart_top = TITLE_TOP + title_height + 0.25  # clear gap below title

    bottom_used = BOTTOM_MARGIN
    if has_footer:
        bottom_used += 0.35
    if has_insight and insight_text:
        insight_lines = measure_text_lines(insight_text, 11, content_width)
        bottom_used += 0.2 + (insight_lines * 0.22)

    chart_height = SLIDE_HEIGHT - chart_top - bottom_used
    chart_height = max(2.5, min(5.5, chart_height))

    return (
        Inches(sidebar_width),
        Inches(chart_top),
        Inches(content_width),
        Inches(chart_height),
        title_font_pt,
    )


def calculate_chart_formatting(num_categories, num_series, chart_height_inches):
    """Calculate font sizes and gap width based on chart density.

    Returns dict: label_size, axis_size, legend_size, gap_width.
    """
    density = num_categories * max(num_series, 1)

    if density <= 10:
        return {'label_size': Pt(12), 'axis_size': Pt(11),
                'legend_size': Pt(10), 'gap_width': 80}
    elif density <= 25:
        return {'label_size': Pt(10), 'axis_size': Pt(10),
                'legend_size': Pt(9), 'gap_width': 50}
    elif density <= 40:
        return {'label_size': Pt(8), 'axis_size': Pt(8),
                'legend_size': Pt(8), 'gap_width': 30}
    else:
        return {'label_size': Pt(7), 'axis_size': Pt(7),
                'legend_size': Pt(7), 'gap_width': 20}


# =============================================================================
# Data helpers
# =============================================================================

def filter_rows(df, exclude=None):
    """Remove rows that shouldn't appear in charts."""
    if exclude is None:
        exclude = DEFAULT_EXCLUDE_ROWS
    else:
        exclude = {r.lower().strip() for r in exclude}
    mask = df.index.str.lower().str.strip().isin(exclude)
    return df[~mask]


def get_series_columns(slide_info, df):
    """Determine which columns to chart based on config and chart type."""
    explicit = slide_info.get('series_columns')
    if explicit:
        if explicit == 'all':
            return list(df.columns)
        elif explicit == 'all_except_total':
            return [c for c in df.columns if c.lower() != 'total']
        elif isinstance(explicit, list):
            return [c for c in explicit if c in df.columns]

    chart_type = slide_info.get('chart_type', 'bar').lower()

    # Stacked charts use all columns except Total
    if 'stacked' in chart_type:
        cols = [c for c in df.columns if c.lower() != 'total']
        return cols if cols else list(df.columns)

    # Clustered with notes hinting at multi-series
    notes = slide_info.get('notes', '').lower()
    if 'clustered' in chart_type or 'gender' in notes or 'region' in notes:
        cols = [c for c in df.columns if c.lower() != 'total']
        if len(cols) > 1:
            return cols

    # Default: single series on Total
    if 'Total' in df.columns:
        return ['Total']
    return [df.columns[0]] if len(df.columns) > 0 else []


def resolve_data_path(data_file, data_dir):
    """Resolve a data file path — absolute or relative to data_dir."""
    if os.path.isabs(data_file):
        return data_file
    return os.path.join(data_dir, data_file)


def parse_color(hex_str):
    """Parse '#RRGGBB' to RGBColor."""
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def normalize_slide_type(raw_type):
    """Normalize any slide type string to a canonical type."""
    return SLIDE_TYPE_ALIASES.get(raw_type.lower().strip(), raw_type.lower().strip())


# =============================================================================
# Chart builder
# =============================================================================

def build_chart(slide, chart_type_str, df, series_columns,
                left, top, width, height, colors, single_color,
                transpose=False):
    """Build a chart with ALL formatting rules applied.

    Returns the chart Shape (for positioning insight/footer relative to it).
    Returns None if the data is empty or chart creation fails.

    transpose: If True, swap rows and columns (PowerPoint "Switch Row/Column").
               For stacked charts this puts demographics on Y-axis and
               answer options as stacking series. Equivalent to the VBA
               Chart.PlotBy toggle Matt uses in his macro.
    """
    if len(df) == 0 or not series_columns:
        return None

    # Clean data
    df = df.apply(pd.to_numeric, errors='coerce').fillna(0)

    xl_type = CHART_TYPE_MAP.get(chart_type_str.lower().strip(),
                                  XL_CHART_TYPE.BAR_CLUSTERED)
    is_bar = xl_type in (XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.BAR_STACKED_100)
    is_stacked = xl_type in (XL_CHART_TYPE.BAR_STACKED_100,
                              XL_CHART_TYPE.COLUMN_STACKED_100)
    is_pie = xl_type == XL_CHART_TYPE.PIE

    # Transpose for "switch rows/columns" — stacked charts typically need this
    # so demographics appear on Y-axis and answer options stack as series
    if transpose or is_stacked:
        # Only use the requested series columns before transposing
        valid_cols = [c for c in series_columns if c in df.columns]
        if valid_cols:
            df = df[valid_cols].T
        else:
            df = df.T
        # After transpose: old columns are now rows (categories),
        # old rows are now columns (series)
        series_columns = list(df.columns)

    # Build chart data
    chart_data = CategoryChartData()
    chart_data.categories = list(df.index)

    valid_cols = [c for c in series_columns if c in df.columns]
    if not valid_cols:
        return None

    for col in valid_cols:
        chart_data.add_series(col, tuple(df[col]))

    num_series = len(valid_cols)

    # Create chart
    chart_shape = slide.shapes.add_chart(xl_type, left, top, width, height, chart_data)
    chart = chart_shape.chart

    # --- Hide chart title (suppresses "Total" label on single-series) ---
    chart.has_title = False

    # --- Height in inches for formatting calc ---
    chart_h_inches = height / 914400 if isinstance(height, int) else 4.0
    fmt = calculate_chart_formatting(len(df), num_series, chart_h_inches)

    # --- Value axis: hide ---
    if not is_pie:
        chart.value_axis.has_major_gridlines = False
        chart.value_axis.visible = False

    # --- Category axis: reverse order on ALL bar types ---
    if is_bar:
        chart.category_axis.reverse_order = True
        chart.category_axis.has_major_gridlines = False

    # --- Category axis font ---
    if not is_pie and hasattr(chart.category_axis, 'tick_labels'):
        chart.category_axis.tick_labels.font.size = fmt['axis_size']

    # --- Legend ---
    if is_pie:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = fmt['legend_size']
    elif num_series > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = fmt['legend_size']
    else:
        chart.has_legend = False

    # --- Plot formatting ---
    plot = chart.plots[0]

    if not is_pie:
        plot.gap_width = fmt['gap_width']

    # Data labels — ALWAYS '0"%"', NEVER '0%'
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    dl.number_format = '0"%"'
    dl.font.size = fmt['label_size']

    if is_pie:
        dl.show_category_name = False
        dl.show_percentage = False

    # --- Colors ---
    color_list = [parse_color(c) for c in colors]
    single_rgb = parse_color(single_color)

    if is_pie:
        for i, point in enumerate(chart.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = color_list[i % len(color_list)]
    elif is_stacked or num_series > 1:
        for i, series in enumerate(chart.series):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color_list[i % len(color_list)]
    else:
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = single_rgb

    return chart_shape


# =============================================================================
# Insight and footer positioning
# =============================================================================

def add_insight_and_footer(slide, chart_shape, insight_text='', footer_text='',
                           sidebar_width=DEFAULT_SIDEBAR_WIDTH):
    """Position insight and footer BELOW the chart, relative to chart bottom.

    This is the single codepath for insight/footer placement.
    Insight is always below the chart, never beside it.
    """
    content_left = sidebar_width
    content_width = SLIDE_WIDTH - sidebar_width - 0.3
    max_bottom = 7.3  # leave margin from slide bottom (7.5")

    chart_bottom = (chart_shape.top + chart_shape.height) / 914400
    y = chart_bottom + 0.08

    if insight_text and y < max_bottom:
        remaining = max_bottom - y
        insight_lines = measure_text_lines(insight_text, 11, content_width)
        insight_height = min(0.7, max(0.3, insight_lines * 0.2))

        if y + insight_height <= max_bottom:
            box = slide.shapes.add_textbox(
                Inches(content_left), Inches(y),
                Inches(content_width), Inches(insight_height),
            )
            tf = box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = insight_text
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.bold = True
            y += insight_height + 0.05

    if footer_text and y < max_bottom:
        footer_y = min(y, max_bottom - 0.25)
        box = slide.shapes.add_textbox(
            Inches(content_left), Inches(footer_y),
            Inches(content_width), Inches(0.25),
        )
        tf = box.text_frame
        tf.paragraphs[0].text = footer_text
        tf.paragraphs[0].font.size = Pt(8)
        tf.paragraphs[0].font.color.rgb = RGBColor(0x80, 0x80, 0x80)


# =============================================================================
# Slide builders
# =============================================================================

def _set_title(slide, text, auto_fit=True):
    """Set the slide title with optional auto-fit."""
    if slide.shapes.title is None:
        return
    slide.shapes.title.text = text
    if auto_fit:
        slide.shapes.title.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _set_footer_placeholder(slide, text, ph_idx):
    """Set footer text in a placeholder, silently skip if not found."""
    try:
        slide.placeholders[ph_idx].text = text
    except (KeyError, IndexError):
        pass


def build_title_slide(prs, info, get_layout):
    slide = prs.slides.add_slide(get_layout('title'))
    _set_title(slide, info.get('title', ''))
    return slide


def _body_font_size(text, available_height_inches=5.1):
    """Choose body font size based on text length to prevent overflow.

    The One column slide body placeholder is ~5.1" tall.
    With widened placeholder (~10") and tight line spacing, more text fits.
    Returns font size in Pt.
    """
    line_count = text.count('\n') + 1
    char_count = len(text)

    # Template body placeholder is 5.51" wide x 5.13" tall
    # At 0.95 line spacing:
    # 12pt ~ 60 chars/line, ~28 lines
    # 11pt ~ 65 chars/line, ~31 lines
    # 10pt ~ 72 chars/line, ~34 lines
    #  9pt ~ 80 chars/line, ~38 lines

    if line_count <= 15 and char_count < 700:
        return Pt(12)
    elif line_count <= 22 and char_count < 1100:
        return Pt(11)
    elif line_count <= 30 and char_count < 1600:
        return Pt(10)
    else:
        return Pt(9)


def build_text_slide(prs, info, get_layout, ph_map, sidebar_width):
    slide = prs.slides.add_slide(get_layout('text'))
    _set_title(slide, info.get('title', ''))

    body = info.get('body', '')
    if body:
        body_text = body if isinstance(body, str) else '\n'.join(body)
        font_size = _body_font_size(body_text)
        body_idx = ph_map.get('body', 17)
        written = False

        # Try the configured placeholder index
        try:
            ph = slide.placeholders[body_idx]
            tf = ph.text_frame

            # Support array-of-paragraphs or newline-separated string
            if isinstance(body, list):
                tf.text = body[0]
                for para_text in body[1:]:
                    p = tf.add_paragraph()
                    p.text = para_text
            else:
                tf.text = body

            # Apply font size and tight line spacing to all paragraphs
            for para in tf.paragraphs:
                para.font.size = font_size
                para.space_after = Pt(2)
                para.space_before = Pt(0)
                para.line_spacing = 0.95

            tf.word_wrap = True
            # Let PowerPoint auto-shrink as a safety net, but our font
            # sizing should already fit
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            written = True
        except (KeyError, IndexError):
            pass

        # Fallback: find body placeholder by type
        if not written:
            for ph in slide.placeholders:
                if ph.placeholder_format.type == 2:  # BODY
                    ph.text_frame.text = body_text
                    for para in ph.text_frame.paragraphs:
                        para.font.size = font_size
                        para.space_after = Pt(2)
                        para.space_before = Pt(0)
                        para.line_spacing = 0.95
                    ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    break

    footer = info.get('footer', '')
    if footer:
        _set_footer_placeholder(slide, footer, ph_map.get('text_footer', 18))

    return slide


def build_section_slide(prs, info, section_counter):
    """Build a section divider with gradient cycling."""
    layout_name = section_counter.next()

    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name.strip() == layout_name.strip():
            slide = prs.slides.add_slide(layout)
            _set_title(slide, info.get('title', ''))
            return slide

    # Should not happen — fall through to first available
    raise ValueError(f"Section layout '{layout_name}' not found in template")


def build_chart_slide(prs, info, get_layout, ph_map, colors, single_color,
                      sidebar_width, data_dir, verbose=False):
    slide = prs.slides.add_slide(get_layout('chart'))

    title = info.get('title', '')
    insight = info.get('insight', '')
    footer = info.get('footer', '')

    _set_title(slide, title)

    # Adaptive layout
    left, top, width, height, title_font_pt = calculate_layout(
        title, sidebar_width,
        has_insight=bool(insight), insight_text=insight,
        has_footer=bool(footer),
    )

    # Set title font size from calculator
    if slide.shapes.title:
        for para in slide.shapes.title.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(title_font_pt)
            if not para.runs:
                para.font.size = Pt(title_font_pt)

    # Resolve data
    data_file = info.get('data_file', '')
    if not data_file:
        print(f'    no data_file, slide without chart')
        _set_footer_placeholder(slide, footer, ph_map.get('graph_footer', 22))
        return slide

    data_path = resolve_data_path(data_file, data_dir)
    if not os.path.exists(data_path):
        print(f'    data not found: {data_path}')
        _set_footer_placeholder(slide, footer, ph_map.get('graph_footer', 22))
        return slide

    try:
        df = to_dataframe(data_path)
    except Exception as e:
        print(f'    data load failed: {e}')
        _set_footer_placeholder(slide, footer, ph_map.get('graph_footer', 22))
        return slide

    # Filter rows
    exclude = info.get('exclude_rows')
    if exclude is None:
        exclude = list(DEFAULT_EXCLUDE_ROWS)
    df = filter_rows(df, exclude)

    if len(df) == 0:
        print(f'    all rows filtered, slide without chart')
        _set_footer_placeholder(slide, footer, ph_map.get('graph_footer', 22))
        return slide

    # Series columns
    series_columns = get_series_columns(info, df)
    if not series_columns:
        print(f'    no valid series columns')
        _set_footer_placeholder(slide, footer, ph_map.get('graph_footer', 22))
        return slide

    chart_type = info.get('chart_type', 'bar')

    if verbose:
        print(f'    data: {len(df)} rows x {len(series_columns)} series ({chart_type})')

    # Build chart — transpose flag from manifest or auto for stacked
    transpose = info.get('transpose', False)
    chart_shape = build_chart(
        slide, chart_type, df, series_columns,
        left, top, width, height,
        colors, single_color,
        transpose=transpose,
    )

    # Insight + footer below chart
    if chart_shape:
        add_insight_and_footer(slide, chart_shape, insight, footer, sidebar_width)
    else:
        print(f'    chart build returned None (empty data?)')
        _set_footer_placeholder(slide, footer, ph_map.get('graph_footer', 22))

    return slide


def build_appendix_slide(prs, info, get_layout):
    slide = prs.slides.add_slide(get_layout('appendix'))
    _set_title(slide, info.get('title', ''))
    return slide


def build_thankyou_slide(prs, info, get_layout):
    slide = prs.slides.add_slide(get_layout('thankyou'))
    _set_title(slide, info.get('title', 'Thank You'))
    return slide


# =============================================================================
# Template helpers
# =============================================================================

class SectionCounter:
    """Cycle through section divider gradient layouts."""

    def __init__(self, gradient_layouts):
        self.layouts = gradient_layouts
        self.idx = 0

    def next(self):
        name = self.layouts[self.idx % len(self.layouts)]
        self.idx += 1
        return name


def make_layout_getter(prs, layout_map):
    """Create a layout lookup function from the template."""
    cache = {}

    def get_layout(slide_type):
        if slide_type in cache:
            return cache[slide_type]

        name = layout_map.get(slide_type)
        if isinstance(name, list):
            name = name[0]

        if not name:
            available = [l.name.strip() for l in prs.slide_masters[0].slide_layouts]
            raise ValueError(
                f"No layout mapping for '{slide_type}'. "
                f"Available layouts: {available}"
            )

        for layout in prs.slide_masters[0].slide_layouts:
            if layout.name.strip() == name.strip():
                cache[slide_type] = layout
                return layout

        available = [l.name.strip() for l in prs.slide_masters[0].slide_layouts]
        raise ValueError(
            f"Layout '{name}' not found in template. "
            f"Available: {available}"
        )

    return get_layout


def discover_layouts(prs):
    """Print all layouts and their placeholders for debugging."""
    for layout in prs.slide_masters[0].slide_layouts:
        phs = []
        for ph in layout.placeholders:
            phs.append(f'idx={ph.placeholder_format.idx} "{ph.name}"')
        print(f'  {layout.name}: {", ".join(phs) or "(no placeholders)"}')


# =============================================================================
# Post-generation validator
# =============================================================================

def validate_pptx(pptx_path, verbose=False):
    """Check a generated PPTX for common issues.

    Returns a list of issue strings. Empty list = clean.
    """
    prs = Presentation(pptx_path)
    issues = []

    for i, slide in enumerate(prs.slides):
        sn = i + 1
        layout = slide.slide_layout.name.strip() if slide.slide_layout else 'unknown'

        # Empty title
        if slide.shapes.title and not slide.shapes.title.text.strip():
            issues.append(f'Slide {sn} ({layout}): empty title')

        charts = [s for s in slide.shapes if s.has_chart]

        for cs in charts:
            chart = cs.chart

            # Zero-dimension chart
            if cs.width == 0 or cs.height == 0:
                issues.append(f'Slide {sn}: chart has 0 dimensions')

            # Chart title showing (should be hidden)
            if chart.has_title:
                issues.append(f'Slide {sn}: chart title visible (should be hidden)')

            # Bar chart without reverse_order
            ct = chart.chart_type
            bar_types = (
                XL_CHART_TYPE.BAR_CLUSTERED,
                XL_CHART_TYPE.BAR_STACKED,
                XL_CHART_TYPE.BAR_STACKED_100,
            )
            if ct in bar_types and not chart.category_axis.reverse_order:
                issues.append(f'Slide {sn}: bar chart missing reverse_order')

            # Legend on single-series non-pie
            if ct != XL_CHART_TYPE.PIE and len(chart.series) == 1 and chart.has_legend:
                issues.append(f'Slide {sn}: single-series chart has legend')

            # Wrong number format
            plot = chart.plots[0]
            if plot.has_data_labels:
                nf = plot.data_labels.number_format
                if nf == '0%':
                    issues.append(f'Slide {sn}: number_format is 0% (should be 0"%")')

        if verbose:
            print(f'  Slide {sn}: {layout} [{len(charts)} chart(s)]')

    return issues


# =============================================================================
# Main generation
# =============================================================================

def generate(args):
    """Generate a PPTX report from manifest + template."""

    # Load manifest
    with open(args.manifest) as f:
        manifest_data = json.load(f)

    if isinstance(manifest_data, list):
        slides = manifest_data
        config = {}
    else:
        slides = manifest_data.get('slides', [])
        config = manifest_data.get('config', {})

    if not slides:
        print('ERROR: No slides in manifest')
        sys.exit(1)

    # Resolve settings: CLI flags > manifest config > defaults
    colors = (
        [c.strip() for c in args.colors.split(',')]
        if args.colors else config.get('colors', DEFAULT_COLORS)
    )
    single_color = args.single_color or config.get('single_color', colors[0])
    sidebar_width = args.sidebar_width or config.get('sidebar_width', DEFAULT_SIDEBAR_WIDTH)
    data_dir = args.data_dir or config.get('data_dir', os.path.dirname(os.path.abspath(args.manifest)))

    layout_map = dict(DEFAULT_LAYOUT_MAP)
    if config.get('layout_map'):
        layout_map.update(config['layout_map'])
    if args.layout_map:
        layout_map.update(json.loads(args.layout_map))

    ph_map = dict(DEFAULT_PLACEHOLDER_MAP)
    if config.get('placeholder_map'):
        ph_map.update(config['placeholder_map'])

    # Load template
    prs = Presentation(args.template)

    if args.list_layouts:
        discover_layouts(prs)
        sys.exit(0)

    get_layout = make_layout_getter(prs, layout_map)

    # Section gradient cycling
    gradients = layout_map.get('section', DEFAULT_LAYOUT_MAP['section'])
    if isinstance(gradients, str):
        gradients = [gradients]
    section_counter = SectionCounter(gradients)

    # Generate
    total = len(slides)
    print(f'Generating {total} slides from {args.template}')

    for i, info in enumerate(slides):
        raw_type = info.get('type', 'chart')
        slide_type = normalize_slide_type(raw_type)
        title = info.get('title', f'Slide {i + 1}')
        label = title[:65] + '...' if len(title) > 65 else title
        print(f'  [{i + 1:>2}/{total}] {slide_type:<8} {label}')

        try:
            if slide_type == 'title':
                build_title_slide(prs, info, get_layout)
            elif slide_type == 'text':
                build_text_slide(prs, info, get_layout, ph_map, sidebar_width)
            elif slide_type == 'section':
                build_section_slide(prs, info, section_counter)
            elif slide_type == 'chart':
                build_chart_slide(prs, info, get_layout, ph_map,
                                  colors, single_color, sidebar_width,
                                  data_dir, verbose=args.verbose)
            elif slide_type == 'appendix':
                build_appendix_slide(prs, info, get_layout)
            elif slide_type == 'thankyou':
                build_thankyou_slide(prs, info, get_layout)
            else:
                print(f'    unknown type "{raw_type}", skipping')
        except Exception as e:
            print(f'    ERROR: {e}')
            if args.verbose:
                import traceback
                traceback.print_exc()

    # Save
    prs.save(args.output)
    print(f'\nSaved: {args.output} ({total} slides)')

    # Post-generation validation
    if args.validate:
        print('\nValidating...')
        issues = validate_pptx(args.output, verbose=args.verbose)
        if issues:
            print(f'{len(issues)} issue(s):')
            for issue in issues:
                print(f'  ! {issue}')
            return issues
        else:
            print('Clean — no issues found')

    # Thumbnails
    if args.thumbnails:
        thumb_script = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'thumbnail.py')
        if os.path.exists(thumb_script):
            prefix = args.output.rsplit('.', 1)[0] + '_qa'
            cols = args.thumbnail_cols or 4
            cmd = f'python3 "{thumb_script}" "{args.output}" "{prefix}" --cols {cols}'
            print(f'\nGenerating thumbnails...')
            os.system(cmd)
        else:
            print(f'\nThumbnail script not found at {thumb_script}')

    return []


# =============================================================================
# CLI
# =============================================================================

def main():
    p = argparse.ArgumentParser(
        description='Generate branded PPTX reports from a manifest + template.',
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )

    # Modes
    p.add_argument('--validate-only', metavar='PPTX',
                   help='Validate an existing PPTX (no generation)')
    p.add_argument('--list-layouts', action='store_true',
                   help='List template layouts and placeholders, then exit')

    # Required for generation
    p.add_argument('--template', help='Client PPTX template file')
    p.add_argument('--manifest', help='Slide manifest JSON file')
    p.add_argument('--output', help='Output PPTX path')

    # Data
    p.add_argument('--data-dir',
                   help='Directory with data JSON files (default: manifest dir)')

    # Colors
    p.add_argument('--colors',
                   help='Comma-separated hex colors for chart series')
    p.add_argument('--single-color',
                   help='Hex color for single-series bar charts (default: first color)')

    # Layout
    p.add_argument('--sidebar-width', type=float, default=None,
                   help=f'Template sidebar width in inches (default: {DEFAULT_SIDEBAR_WIDTH})')
    p.add_argument('--layout-map',
                   help='JSON string overriding layout names per slide type')

    # Post-generation
    p.add_argument('--validate', action='store_true',
                   help='Run validation checks after generation')
    p.add_argument('--thumbnails', action='store_true',
                   help='Generate QA thumbnail grids after generation')
    p.add_argument('--thumbnail-cols', type=int, default=4,
                   help='Columns per thumbnail grid (default: 4)')

    # Debug
    p.add_argument('-v', '--verbose', action='store_true',
                   help='Verbose output with data details')
    p.add_argument('--dry-run', action='store_true',
                   help='Parse manifest and show plan without generating')

    args = p.parse_args()

    # Validate-only mode
    if args.validate_only:
        issues = validate_pptx(args.validate_only, verbose=True)
        if issues:
            print(f'\n{len(issues)} issue(s):')
            for issue in issues:
                print(f'  ! {issue}')
            sys.exit(1)
        else:
            print('Clean — no issues found')
            sys.exit(0)

    # List layouts mode (needs --template)
    if args.list_layouts:
        if not args.template:
            p.error('--list-layouts requires --template')
        prs = Presentation(args.template)
        discover_layouts(prs)
        sys.exit(0)

    # Generation mode
    if not all([args.template, args.manifest, args.output]):
        p.error('--template, --manifest, and --output are required')

    if not os.path.exists(args.template):
        p.error(f'Template not found: {args.template}')
    if not os.path.exists(args.manifest):
        p.error(f'Manifest not found: {args.manifest}')

    # Dry run
    if args.dry_run:
        with open(args.manifest) as f:
            data = json.load(f)
        slides = data if isinstance(data, list) else data.get('slides', [])
        print(f'Manifest: {len(slides)} slides')
        for i, s in enumerate(slides):
            st = normalize_slide_type(s.get('type', 'chart'))
            t = s.get('title', 'untitled')
            label = t[:70] + '...' if len(t) > 70 else t
            extras = []
            if s.get('data_file'):
                extras.append(f'data={s["data_file"]}')
            if s.get('chart_type'):
                extras.append(f'chart={s["chart_type"]}')
            if s.get('insight'):
                extras.append('has insight')
            suffix = f' ({", ".join(extras)})' if extras else ''
            print(f'  [{i + 1:>2}] {st:<8} {label}{suffix}')
        sys.exit(0)

    issues = generate(args)
    sys.exit(1 if issues else 0)


if __name__ == '__main__':
    main()
