"""
PowerPoint report template library.

Generates native-editable PPTX reports with charts, KPIs, and text slides.
Supports two modes: blank (default Platinum style) or template-based (user's
reference PPTX preserving all design tokens).

Usage without reference template (default Platinum style):
    from pptx_template import ReportTemplate
    report = ReportTemplate()
    report.set_palette({'primary': '#0072CE', 'series': ['#0072CE', '#6400AA', '#C8102E']})
    report.title_slide('Brand Health Report', 'Q1 2026')
    report.chart_slide('Aided Awareness', df, chart_type='bar')
    report.save()

Usage with reference template (preserves all design):
    report = ReportTemplate(template_path='/workspace/uploads/client_template.pptx')
    print(json.dumps(report.design, indent=2))
    report.title_slide('Brand Health Report', 'Q1 2026')
    report.save()
"""

import json
import os
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree


# Default series palette — Platinum brand
_DEFAULT_SERIES = [
    "#3b82f6",  # blue
    "#ef4444",  # red
    "#22c55e",  # green
    "#f59e0b",  # amber
    "#8b5cf6",  # violet
    "#ec4899",  # pink
    "#06b6d4",  # cyan
    "#f97316",  # orange
]

_CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}

# XML namespace map for theme parsing
_NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}


class ReportTemplate:
    """PowerPoint report builder with optional template support.

    Args:
        template_path: Path to a reference .pptx file. If provided, the
            template's theme, colors, fonts, and layouts are preserved.
            If None, a blank presentation with Platinum brand defaults is
            created.
        widescreen: If True (default), use 16:9 aspect ratio (13.333" x 7.5").
            If False, use 4:3 (10" x 7.5"). Ignored when template_path is set.
    """

    def __init__(self, template_path=None, widescreen=True):
        self._is_dark = False
        self._dark_bg = "#0f172a"
        self._dark_text = "#FFFFFF"

        if template_path and os.path.isfile(template_path):
            self.prs = Presentation(template_path)
            self._strip_content_slides()
            self.design = self._extract_design_spec()
        else:
            self.prs = Presentation()
            if widescreen:
                self.prs.slide_width = Inches(13.333)
                self.prs.slide_height = Inches(7.5)
            else:
                self.prs.slide_width = Inches(10)
                self.prs.slide_height = Inches(7.5)
            self.design = self._default_design()

        self._palette = {
            "primary": _DEFAULT_SERIES[0],
            "series": list(_DEFAULT_SERIES),
        }

    # ------------------------------------------------------------------
    # Internal: template handling
    # ------------------------------------------------------------------

    def _strip_content_slides(self):
        """Remove all content slides, keeping slide masters and layouts."""
        sldIdLst = self.prs.slides._sldIdLst
        rIds = []
        for sldId in list(sldIdLst):
            rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
            if rId:
                rIds.append((sldId, rId))
        for sldId, rId in rIds:
            self.prs.part.drop_rel(rId)
            sldIdLst.remove(sldId)

    def _extract_design_spec(self):
        """Parse theme XML from the first slide master and return design spec.

        Returns:
            Dict with 'colors', 'fonts', and 'layouts' keys.
        """
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT

        colors = {}
        fonts = {"heading": "Calibri", "body": "Calibri"}
        layouts = []

        try:
            theme_part = self.prs.slide_masters[0].part.part_related_by(RT.THEME)
            theme_xml = theme_part.blob
            root = etree.fromstring(theme_xml)

            # --- Color scheme ---
            clr_scheme = root.find(".//a:clrScheme", _NSMAP)
            if clr_scheme is not None:
                slot_names = [
                    "dk1", "lt1", "dk2", "lt2",
                    "accent1", "accent2", "accent3", "accent4",
                    "accent5", "accent6", "hlink", "folHlink",
                ]
                for name in slot_names:
                    el = clr_scheme.find(f"a:{name}", _NSMAP)
                    if el is not None:
                        srgb = el.find("a:srgbClr", _NSMAP)
                        sys_clr = el.find("a:sysClr", _NSMAP)
                        if srgb is not None:
                            colors[name] = "#" + srgb.get("val", "000000")
                        elif sys_clr is not None:
                            colors[name] = "#" + sys_clr.get("lastClr", "000000")

            # --- Font scheme ---
            font_scheme = root.find(".//a:fontScheme", _NSMAP)
            if font_scheme is not None:
                major = font_scheme.find("a:majorFont/a:latin", _NSMAP)
                minor = font_scheme.find("a:minorFont/a:latin", _NSMAP)
                if major is not None:
                    fonts["heading"] = major.get("typeface", "Calibri")
                if minor is not None:
                    fonts["body"] = minor.get("typeface", "Calibri")
        except Exception:
            pass

        # --- Layouts ---
        for master in self.prs.slide_masters:
            for layout in master.slide_layouts:
                ph_info = []
                for ph in layout.placeholders:
                    ph_info.append({
                        "idx": ph.placeholder_format.idx,
                        "name": ph.name,
                        "type": str(ph.placeholder_format.type),
                    })
                layouts.append({
                    "name": layout.name,
                    "placeholders": ph_info,
                })

        return {"colors": colors, "fonts": fonts, "layouts": layouts}

    def _default_design(self):
        """Return a default design spec with Platinum brand colors."""
        return {
            "colors": {
                "dk1": "#0f172a",
                "lt1": "#f1f5f9",
                "dk2": "#1e293b",
                "lt2": "#e2e8f0",
                "accent1": "#3b82f6",
                "accent2": "#ef4444",
                "accent3": "#22c55e",
                "accent4": "#f59e0b",
                "accent5": "#8b5cf6",
                "accent6": "#ec4899",
                "hlink": "#3b82f6",
                "folHlink": "#8b5cf6",
            },
            "fonts": {
                "heading": "Calibri",
                "body": "Calibri",
            },
            "layouts": [],
        }

    # ------------------------------------------------------------------
    # Layout helpers
    # ------------------------------------------------------------------

    def get_layout_by_name(self, name):
        """Find a slide layout by fuzzy name match.

        Args:
            name: Layout name to search for (case-insensitive substring match).

        Returns:
            The first matching SlideLayout, or the first available layout as
            fallback.
        """
        target = name.lower()
        for master in self.prs.slide_masters:
            for layout in master.slide_layouts:
                if target in layout.name.lower():
                    return layout
        # Fallback to first layout
        return self.prs.slide_masters[0].slide_layouts[0]

    # ------------------------------------------------------------------
    # Palette and theme
    # ------------------------------------------------------------------

    def set_palette(self, palette):
        """Set the color palette for charts and accents.

        Args:
            palette: Dict with 'primary' (hex string) and 'series' (list of
                hex strings) keys.
        """
        self._palette = palette

    def set_dark_theme(self, bg_hex="#0f172a", text_hex="#FFFFFF"):
        """Enable dark slide backgrounds.

        Args:
            bg_hex: Background color as hex string.
            text_hex: Text color as hex string.
        """
        self._dark_bg = bg_hex
        self._dark_text = text_hex
        self._is_dark = True

    # ------------------------------------------------------------------
    # Slide builders
    # ------------------------------------------------------------------

    def title_slide(self, title, subtitle="", date=""):
        """Add a title slide.

        Args:
            title: Main title text.
            subtitle: Optional subtitle line.
            date: Optional date string (appended to subtitle).
        """
        layout = self.get_layout_by_name("Title")
        slide = self.prs.slides.add_slide(layout)

        if self._is_dark:
            self._set_slide_bg(slide, self._dark_bg)

        # Set title placeholder
        if slide.placeholders:
            slide.placeholders[0].text = title
            if self._is_dark:
                for para in slide.placeholders[0].text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = self._hex_to_rgb(self._dark_text)

        # Subtitle (placeholder index 1 if it exists)
        sub_text = subtitle
        if date:
            sub_text = f"{subtitle}\n{date}" if subtitle else date
        if sub_text and len(slide.placeholders) > 1:
            slide.placeholders[1].text = sub_text
            if self._is_dark:
                for para in slide.placeholders[1].text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = self._hex_to_rgb(self._dark_text)

        return slide

    def chart_slide(self, title, df, chart_type="bar", narrative="", source=""):
        """Add a slide with a native editable chart.

        Args:
            title: Slide title text.
            df: pandas DataFrame. Index = categories, columns = series names.
            chart_type: One of 'bar', 'column', 'line', 'pie', 'doughnut'.
            narrative: Optional text block displayed to the right of the chart.
            source: Optional source attribution at the bottom of the slide.

        Returns:
            The created slide object.
        """
        layout = self.get_layout_by_name("Blank")
        slide = self.prs.slides.add_slide(layout)

        if self._is_dark:
            self._set_slide_bg(slide, self._dark_bg)

        slide_w = self.prs.slide_width
        slide_h = self.prs.slide_height

        # Title text box
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), slide_w - Inches(1), Inches(0.6),
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.name = self.design["fonts"]["heading"]
        if self._is_dark:
            p.font.color.rgb = self._hex_to_rgb(self._dark_text)

        # Chart data
        chart_data = CategoryChartData()
        chart_data.categories = list(df.index)
        for col in df.columns:
            chart_data.add_series(str(col), list(df[col]))

        # Chart position: left 60% if narrative, otherwise full width
        chart_left = Inches(0.5)
        chart_top = Inches(1.2)
        if narrative:
            chart_width = int(slide_w * 0.55)
        else:
            chart_width = slide_w - Inches(1)
        chart_height = slide_h - Inches(2.2)

        xl_type = _CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.BAR_CLUSTERED)
        chart_frame = slide.shapes.add_chart(
            xl_type, chart_left, chart_top, chart_width, chart_height, chart_data,
        )
        chart = chart_frame.chart
        self._apply_chart_palette(chart)

        if self._is_dark:
            self._style_chart_dark(chart)

        # Narrative text box on right
        if narrative:
            nar_left = int(slide_w * 0.60)
            nar_width = slide_w - nar_left - Inches(0.5)
            nar_box = slide.shapes.add_textbox(
                nar_left, Inches(1.2), nar_width, slide_h - Inches(2.5),
            )
            ntf = nar_box.text_frame
            ntf.word_wrap = True
            np_ = ntf.paragraphs[0]
            np_.text = narrative
            np_.font.size = Pt(12)
            np_.font.name = self.design["fonts"]["body"]
            if self._is_dark:
                np_.font.color.rgb = self._hex_to_rgb(self._dark_text)

        # Source attribution
        if source:
            src_box = slide.shapes.add_textbox(
                Inches(0.5), slide_h - Inches(0.5), slide_w - Inches(1), Inches(0.4),
            )
            stf = src_box.text_frame
            sp = stf.paragraphs[0]
            sp.text = source
            sp.font.size = Pt(8)
            sp.font.italic = True
            sp.font.name = self.design["fonts"]["body"]
            sp.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        return slide

    def kpi_slide(self, title, kpis):
        """Add a slide with KPI metric cards.

        Args:
            title: Slide title text.
            kpis: List of dicts, each with 'label', 'value', and optional
                'change' keys. Example:
                [{'label': 'NPS', 'value': '+42', 'change': '+3 vs Q4'}]
        """
        layout = self.get_layout_by_name("Blank")
        slide = self.prs.slides.add_slide(layout)

        if self._is_dark:
            self._set_slide_bg(slide, self._dark_bg)

        slide_w = self.prs.slide_width
        slide_h = self.prs.slide_height

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), slide_w - Inches(1), Inches(0.6),
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.name = self.design["fonts"]["heading"]
        if self._is_dark:
            p.font.color.rgb = self._hex_to_rgb(self._dark_text)

        # KPI cards
        n = max(len(kpis), 1)
        margin = Inches(0.5)
        gap = Inches(0.3)
        total_gap = gap * (n - 1)
        card_w = int((slide_w - 2 * margin - total_gap) / n)
        card_h = Inches(2.5)
        card_top = Inches(2.5)

        primary_rgb = self._hex_to_rgb(
            self._palette.get("primary", _DEFAULT_SERIES[0])
        )

        for i, kpi in enumerate(kpis):
            left = int(margin + i * (card_w + gap))

            # Rounded rectangle card
            from pptx.enum.shapes import MSO_SHAPE
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, card_top, card_w, card_h,
            )
            card.fill.solid()
            if self._is_dark:
                card.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
            else:
                card.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
            card.line.fill.background()

            # Value
            val_box = slide.shapes.add_textbox(
                left + Inches(0.2), card_top + Inches(0.3),
                card_w - Inches(0.4), Inches(1.0),
            )
            vp = val_box.text_frame.paragraphs[0]
            vp.text = str(kpi.get("value", ""))
            vp.font.size = Pt(36)
            vp.font.bold = True
            vp.font.color.rgb = primary_rgb
            vp.font.name = self.design["fonts"]["heading"]
            vp.alignment = PP_ALIGN.CENTER

            # Label
            lbl_box = slide.shapes.add_textbox(
                left + Inches(0.2), card_top + Inches(1.3),
                card_w - Inches(0.4), Inches(0.6),
            )
            lp = lbl_box.text_frame.paragraphs[0]
            lp.text = str(kpi.get("label", ""))
            lp.font.size = Pt(14)
            lp.font.name = self.design["fonts"]["body"]
            lp.alignment = PP_ALIGN.CENTER
            if self._is_dark:
                lp.font.color.rgb = self._hex_to_rgb(self._dark_text)

            # Change indicator
            change = kpi.get("change", "")
            if change:
                chg_box = slide.shapes.add_textbox(
                    left + Inches(0.2), card_top + Inches(1.8),
                    card_w - Inches(0.4), Inches(0.4),
                )
                cp = chg_box.text_frame.paragraphs[0]
                cp.text = str(change)
                cp.font.size = Pt(11)
                cp.font.name = self.design["fonts"]["body"]
                cp.alignment = PP_ALIGN.CENTER
                # Green for positive, red for negative
                if str(change).startswith("+"):
                    cp.font.color.rgb = RGBColor(0x22, 0xC5, 0x5E)
                elif str(change).startswith("-"):
                    cp.font.color.rgb = RGBColor(0xEF, 0x44, 0x44)
                else:
                    cp.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        return slide

    def text_slide(self, title, bullets):
        """Add a slide with bulleted text.

        Args:
            title: Slide title text.
            bullets: List of strings for bullet points.
        """
        layout = self.get_layout_by_name("Blank")
        slide = self.prs.slides.add_slide(layout)

        if self._is_dark:
            self._set_slide_bg(slide, self._dark_bg)

        slide_w = self.prs.slide_width
        slide_h = self.prs.slide_height

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), slide_w - Inches(1), Inches(0.6),
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.name = self.design["fonts"]["heading"]
        if self._is_dark:
            p.font.color.rgb = self._hex_to_rgb(self._dark_text)

        # Bullets
        body_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.3), slide_w - Inches(1.4), slide_h - Inches(2),
        )
        tf = body_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text = bullet
            para.font.size = Pt(16)
            para.font.name = self.design["fonts"]["body"]
            para.space_after = Pt(8)
            para.level = 0
            if self._is_dark:
                para.font.color.rgb = self._hex_to_rgb(self._dark_text)

        return slide

    def divider_slide(self, section_title):
        """Add a section divider slide with large centered text.

        Args:
            section_title: The section heading text.
        """
        layout = self.get_layout_by_name("Section")
        slide = self.prs.slides.add_slide(layout)

        if self._is_dark:
            self._set_slide_bg(slide, self._dark_bg)

        slide_w = self.prs.slide_width
        slide_h = self.prs.slide_height

        # If the layout has a title placeholder, use it
        if slide.placeholders:
            slide.placeholders[0].text = section_title
            for para in slide.placeholders[0].text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(36)
                    if self._is_dark:
                        run.font.color.rgb = self._hex_to_rgb(self._dark_text)
        else:
            # Blank layout fallback: centered text box
            box = slide.shapes.add_textbox(
                Inches(1), int(slide_h / 2) - Inches(0.5),
                slide_w - Inches(2), Inches(1),
            )
            p = box.text_frame.paragraphs[0]
            p.text = section_title
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.name = self.design["fonts"]["heading"]
            p.alignment = PP_ALIGN.CENTER
            if self._is_dark:
                p.font.color.rgb = self._hex_to_rgb(self._dark_text)

        return slide

    # ------------------------------------------------------------------
    # Save
    # ------------------------------------------------------------------

    def save(self, path="/workspace/report.pptx", max_slides_per_file=15):
        """Save the presentation to disk.

        If the presentation exceeds *max_slides_per_file* slides, it is split
        into multiple part files (report_part1.pptx, report_part2.pptx, ...).

        Args:
            path: Output file path.
            max_slides_per_file: Maximum slides per file before splitting.

        Returns:
            List of file paths written.
        """
        total = len(self.prs.slides)
        if total <= max_slides_per_file:
            os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
            self.prs.save(path)
            return [path]

        # Split into parts
        base, ext = os.path.splitext(path)
        parts = []
        slide_xmls = list(self.prs.slides._sldIdLst)

        part_num = 0
        for start in range(0, total, max_slides_per_file):
            part_num += 1
            chunk_end = min(start + max_slides_per_file, total)

            # Deep-copy approach: save full, then remove unwanted slides
            part_prs = deepcopy(self.prs)
            sldIdLst = part_prs.slides._sldIdLst
            all_ids = list(sldIdLst)

            # Remove slides outside the range
            to_remove = all_ids[:start] + all_ids[chunk_end:]
            for sldId in to_remove:
                rId = sldId.get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
                )
                if rId:
                    try:
                        part_prs.part.drop_rel(rId)
                    except Exception:
                        pass
                sldIdLst.remove(sldId)

            part_path = f"{base}_part{part_num}{ext}"
            os.makedirs(os.path.dirname(part_path) or ".", exist_ok=True)
            part_prs.save(part_path)
            parts.append(part_path)

        return parts

    def save_design_spec(self, path="/workspace/design-spec.json"):
        """Write the extracted design spec to a JSON file.

        Args:
            path: Output file path.

        Returns:
            The path written.
        """
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        with open(path, "w") as f:
            json.dump(self.design, f, indent=2)
        return path

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _hex_to_rgb(self, hex_color):
        """Convert a hex color string to an RGBColor.

        Args:
            hex_color: Color string like '#3b82f6' or '3b82f6'.

        Returns:
            pptx.dml.color.RGBColor instance.
        """
        h = hex_color.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def _apply_chart_palette(self, chart):
        """Apply series colors from self._palette to chart series.

        Args:
            chart: A python-pptx Chart object.
        """
        series_colors = self._palette.get("series", _DEFAULT_SERIES)
        plot = chart.plots[0]
        for i, series in enumerate(plot.series):
            color_hex = series_colors[i % len(series_colors)]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = self._hex_to_rgb(color_hex)

    def _style_chart_dark(self, chart):
        """Apply dark-theme styling to chart axes and background.

        Args:
            chart: A python-pptx Chart object.
        """
        text_rgb = self._hex_to_rgb(self._dark_text)

        # Style axes
        for axis_attr in ("category_axis", "value_axis"):
            axis = getattr(chart, axis_attr, None)
            if axis is None:
                continue
            try:
                axis.tick_labels.font.color.rgb = text_rgb
                axis.tick_labels.font.size = Pt(9)
            except Exception:
                pass
            try:
                axis.major_gridlines.format.line.fill.background()
            except Exception:
                pass

        # Chart background transparent
        try:
            chart.chart_style = 2
        except Exception:
            pass

    def _set_slide_bg(self, slide, hex_color):
        """Set slide background to a solid color.

        Args:
            slide: A python-pptx Slide object.
            hex_color: Background color as hex string.
        """
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(hex_color)
